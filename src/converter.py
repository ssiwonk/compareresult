"""
Slide and PDF converter module.
Extracts slide texts and renders high-resolution slide images using PyMuPDF and MS PowerPoint COM.
"""
import os
import sys
import tempfile
import pymupdf
from pptx import Presentation
from PIL import Image

def extract_pdf(pdf_path, output_img_dir, dpi=160):
    """
    Renders PDF pages as high-resolution images and extracts text.
    """
    os.makedirs(output_img_dir, exist_ok=True)
    doc = pymupdf.open(pdf_path)
    slides = []
    
    for i, page in enumerate(doc):
        slide_num = i + 1
        img_filename = f"slide_{slide_num:03d}.webp"
        img_path = os.path.join(output_img_dir, img_filename)
        
        # Render high-resolution pixmap
        pix = page.get_pixmap(dpi=dpi)
        # Convert to PIL Image to save as optimized WebP
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.save(img_path, "WEBP", quality=85)
        
        # Extract text & lines
        text = page.get_text().strip()
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        title = lines[0] if lines else f"Slide {slide_num}"
        
        slides.append({
            "num": slide_num,
            "title": title,
            "text": text,
            "lines": lines,
            "image": f"images/{os.path.basename(output_img_dir)}/{img_filename}",
            "width": pix.width,
            "height": pix.height
        })
        
    doc.close()
    return slides

def extract_pptx(pptx_path, output_img_dir, dpi=160):
    """
    Converts PPTX to PDF via MS PowerPoint COM for 100% native rendering,
    then renders high-resolution images and extracts text with python-pptx.
    """
    os.makedirs(output_img_dir, exist_ok=True)
    
    # 1. Extract text structure using python-pptx
    prs = Presentation(pptx_path)
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
        slide_texts.append(texts)
        
    # 2. Render slide images via PowerPoint COM -> PDF -> PyMuPDF
    abs_pptx = os.path.abspath(pptx_path)
    temp_dir = tempfile.mkdtemp()
    temp_pdf = os.path.join(temp_dir, "temp_render.pdf")
    
    ppt_app = None
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        prs_com = ppt_app.Presentations.Open(abs_pptx, WithWindow=False)
        prs_com.SaveAs(temp_pdf, 32)  # 32 = ppSaveAsPDF
        prs_com.Close()
    except Exception as e:
        print(f"  [Warning] PowerPoint COM rendering failed ({e}). Attempting direct conversion.")
        raise e
    finally:
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
                
    # Render PDF pages to images
    doc = pymupdf.open(temp_pdf)
    slides = []
    
    for i, page in enumerate(doc):
        slide_num = i + 1
        img_filename = f"slide_{slide_num:03d}.webp"
        img_path = os.path.join(output_img_dir, img_filename)
        
        pix = page.get_pixmap(dpi=dpi)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        img.save(img_path, "WEBP", quality=85)
        
        # Extract text & lines directly from rendered page for 100% fidelity
        page_text = page.get_text().strip()
        pdf_lines = [line.strip() for line in page_text.split("\n") if line.strip()]
        lines = pdf_lines if pdf_lines else (slide_texts[i] if i < len(slide_texts) else [])
        text = "\n".join(lines)
        title = lines[0] if lines else f"Slide {slide_num}"
        
        slides.append({
            "num": slide_num,
            "title": title,
            "text": text,
            "lines": lines,
            "image": f"images/{os.path.basename(output_img_dir)}/{img_filename}",
            "width": pix.width,
            "height": pix.height
        })
        
    doc.close()
    try:
        os.remove(temp_pdf)
        os.rmdir(temp_dir)
    except Exception:
        pass
        
    return slides

def process_file(file_path, base_output_dir):
    """
    Dispatches to appropriate extractor based on extension.
    """
    fname = os.path.basename(file_path)
    file_id = os.path.splitext(fname)[0]
    safe_folder_name = "".join([c if c.isalnum() or c in ("-", "_") else "_" for c in file_id])
    output_img_dir = os.path.join(base_output_dir, "images", safe_folder_name)
    
    ext = os.path.splitext(fname)[1].lower()
    print(f"[*] Processing {fname} ({ext})...")
    if ext == ".pdf":
        slides = extract_pdf(file_path, output_img_dir)
    elif ext in (".pptx", ".ppt"):
        slides = extract_pptx(file_path, output_img_dir)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
        
    print(f"    -> Extracted {len(slides)} slides to {output_img_dir}")
    return {
        "file_name": fname,
        "folder_name": safe_folder_name,
        "slide_count": len(slides),
        "slides": slides
    }
